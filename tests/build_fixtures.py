from __future__ import annotations

from pathlib import Path


ROOT = Path(__file__).resolve().parent / "fixtures"


def build() -> None:
    ROOT.mkdir(parents=True, exist_ok=True)
    (ROOT / "sample.txt").write_text(
        "Unlocking Transformative Excellence\n\n"
        "Revenue increased 14% to EUR 8.2 million after the Berlin team changed pricing.\n\n"
        "It is important to note that the team utilizes weekly reviews.\n",
        encoding="utf-8",
    )
    (ROOT / "sample.md").write_text(
        "# Innovation and Impact\n\n- We reduced processing time from 9 days to 4 days.\n"
        "- Strategic alignment for future growth\n\n`do_not_rewrite()`\n",
        encoding="utf-8",
    )

    from docx import Document
    doc = Document()
    doc.add_heading("Transformative Excellence", level=0)
    doc.add_paragraph("It is important to note that the team utilizes weekly reviews.")
    doc.add_paragraph("Customer churn fell from 12% to 8% in Q2.", style="List Bullet")
    mixed = doc.add_paragraph()
    mixed.add_run("Strategic ").bold = True
    mixed.add_run("alignment for future growth")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "Owner"
    table.cell(0, 1).text = "Decision"
    table.cell(1, 0).text = "Maria"
    table.cell(1, 1).text = "Launch on 15 August"
    doc.sections[0].header.paragraphs[0].text = "Confidential"
    doc.save(ROOT / "sample.docx")

    from pptx import Presentation
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.chart.data import ChartData
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Transformative Excellence"
    frame = slide.placeholders[1].text_frame
    frame.clear()
    frame.paragraphs[0].text = "It is important to note that the team utilizes weekly reviews."
    p = frame.add_paragraph(); p.text = "Cycle time fell from 9 days to 4 days."; p.level = 0
    box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(5), Inches(0.5))
    box.text_frame.paragraphs[0].add_run().text = "Strategic "
    box.text_frame.paragraphs[0].runs[0].font.bold = True
    box.text_frame.paragraphs[0].add_run().text = "alignment for future growth"
    table_shape = slide.shapes.add_table(2, 2, Inches(7), Inches(4.5), Inches(2.5), Inches(1.2))
    table_shape.table.cell(0, 0).text = "Owner"; table_shape.table.cell(0, 1).text = "Decision"
    table_shape.table.cell(1, 0).text = "Maria"; table_shape.table.cell(1, 1).text = "Launch 15 August"
    data = ChartData(); data.categories = ["Q1", "Q2"]; data.add_series("Revenue", (4.0, 5.0))
    slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(10.2), Inches(4.5), Inches(2.8), Inches(2.2), data)
    slide.notes_slide.notes_text_frame.text = "Speaker note: verify the 14% figure."
    prs.save(ROOT / "sample.pptx")

    from pypdf import PdfWriter
    writer = PdfWriter(); writer.add_blank_page(width=612, height=792); writer.add_metadata({"/Subject": "Synthetic DeSlopAI fixture"})
    with (ROOT / "sample.pdf").open("wb") as handle:
        writer.write(handle)


if __name__ == "__main__":
    build()
