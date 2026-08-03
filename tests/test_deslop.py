from __future__ import annotations

import json
import subprocess
import sys
import tempfile
import unittest
from pathlib import Path


ROOT = Path(__file__).resolve().parents[1]
SKILL = ROOT / "skills" / "deslop-ai"
sys.path.insert(0, str(SKILL / "scripts"))

from build_fixtures import build
from deslop_core import Analyzer, DeslopError, PRIVATE_ROOT, editorial_vector, protected_equal, sha256_file, strict_json_load, validate_request
from formats import extract_source


def request(input_data: dict, operation: str = "audit") -> dict:
    return {
        "schemaVersion": "deslop-request/v1", "id": "synthetic-test", "prompt": "Audit this synthetic fixture.",
        "input": input_data, "operation": operation, "genre": "consulting", "profileId": "none",
        "communicationJob": {"audience": "test reviewers", "outcome": "verify behavior", "takeaway": "The tool preserves facts and rejects empty words."},
        "rewritePolicy": "conservative", "preserveTerms": ["EUR"],
    }


class TestDeSlop(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        cls.fixtures = ROOT / "tests" / "fixtures"
        required = [cls.fixtures / name for name in ["sample.txt", "sample.md", "sample.docx", "sample.pptx", "sample.pdf"]]
        if not all(path.exists() for path in required):
            build()
        cls.analyzer = Analyzer()

    def test_catalog_positive_cases_at_least_50(self) -> None:
        terms = []
        for group in self.analyzer.lex_groups:
            terms.extend(str(term) for term in group.get("terms", []))
        self.assertGreaterEqual(len(terms), 50)
        for index, term in enumerate(terms[:50]):
            text = f"This passage mentions {term} in a generic way."
            block = {"blockId": f"b{index}", "locator": f"test:{index}", "text": text, "sourceHash": "h", "scope": "test", "role": "paragraph"}
            self.assertTrue(self.analyzer.analyze_block(block, "general"), term)

    def test_semantic_review_requires_standalone_container_meaning(self) -> None:
        guidance = (SKILL / "references" / "semantic-review.md").read_text(encoding="utf-8")
        skill = (SKILL / "SKILL.md").read_text(encoding="utf-8")
        self.assertIn("Container and standalone-meaning review", guidance)
        self.assertIn("presenter narration", guidance)
        self.assertIn("coherent standalone proposition", guidance)
        self.assertIn("Actionable-title gate", guidance)
        self.assertIn("Recipient-burden review", guidance)
        self.assertIn("source-integrity problem", guidance)
        self.assertIn("definition of done", guidance)
        self.assertIn("What should happen next?", guidance)
        self.assertIn("gating condition", skill)
        self.assertIn("standalone-meaning gate", skill)

    def test_negative_controls_do_not_trigger_high_severity_value_accusation(self) -> None:
        controls = [
            "Revenue increased 14% after the Berlin team changed pricing.",
            "The board rejected Option B because it costs EUR 2 million more.",
            "Smith (2024) reports a 6-point improvement in the treatment group.",
            "We launched the trial on 12 March and enrolled 81 patients.",
            "Cycle time fell from 9 days to 4 days.",
        ] * 5
        blocks = [{"blockId": f"n{i}", "locator": f"negative:{i}", "text": text, "sourceHash": str(i), "scope": f"s{i}", "role": "paragraph"} for i, text in enumerate(controls)]
        _, assessments = self.analyzer.analyze(blocks, "general")
        self.assertEqual(25, len(assessments))
        self.assertFalse(any(item["verdict"] == "needs-improvement" for item in assessments))

    def test_local_fallback_flags_cross_genre_slop_but_keeps_structural_controls(self) -> None:
        slop = [
            "A quiet shift is happening. In today's fast-paced landscape, leaders must embrace a holistic, human-centric mindset. It's not about adopting AI—it's about unlocking potential. The future belongs to those who navigate change with intention.",
            "Leadership isn't about having all the answers. It's about empowering others, fostering trust, and creating space for growth. Let that sink in.",
            "Accelerating strategic transformation through integrated capabilities to unlock sustainable value.",
            "As we navigate this evolving landscape, please leverage cross-functional synergies to drive alignment and operational excellence.",
            "We are thrilled to announce a groundbreaking solution that seamlessly revolutionizes how teams collaborate.",
            "This pivotal study delves into the intricate interplay between X and Y, underscoring the multifaceted nature of the rapidly evolving landscape.",
            "Building a future-ready operating model to enable scalable impact across the enterprise.",
            "Sometimes the biggest lessons come from the smallest moments. Growth isn't linear—it's a journey. Keep showing up.",
            "Experts agree that AI is transforming every industry, and research shows organizations must act now.",
            "Every founder needs three things: clarity, courage, and consistency.",
            "Success isn't about working harder. It's about working smarter.",
        ]
        controls = [
            "We cut onboarding from 14 days to 8 by removing two approval steps.",
            "The board delayed the launch—supplier certification was still incomplete.",
            "The test measured cost, speed, and failure rate.",
            "We estimate treatment effects using a difference-in-differences design across 42 hospitals.",
        ]
        blocks = [
            {"blockId": f"stress-{i}", "locator": f"stress:{i}", "text": text, "sourceHash": str(i), "scope": f"stress-{i}", "role": "paragraph"}
            for i, text in enumerate(slop + controls)
        ]
        _, assessments = self.analyzer.analyze(blocks, "general")
        self.assertTrue(all(item["verdict"] == "needs-improvement" for item in assessments[:len(slop)]))
        self.assertTrue(all(item["verdict"] == "meaningful" for item in assessments[len(slop):]))

    def test_local_fallback_flags_recipient_burden_workslop(self) -> None:
        workslop = [
            ("Next steps", "headline"),
            ("Overview", "headline"),
            ("Please analyze the market and provide strategic recommendations.", "paragraph"),
            ("This document provides a comprehensive overview of key considerations and actionable insights for moving forward.", "paragraph"),
            ("Communication is key. We need to move forward with urgency.", "paragraph"),
            ("The initiative presents both opportunities and challenges. A balanced approach will be essential to ensure success.", "paragraph"),
            ("Continue exploring potential opportunities and engage key stakeholders as appropriate.", "bullet"),
            ("Overall, the findings highlight the importance of adopting a proactive and strategic approach.", "paragraph"),
            ("You asked us to assess the market. The following analysis examines the market and identifies key insights.", "paragraph"),
            ("Organizations should prioritize agility, collaboration, and innovation to remain competitive.", "paragraph"),
        ]
        controls = [
            ("Submit the application by 31 March", "headline"),
            ("Maya will send the revised cost model to Finance by Friday.", "bullet"),
            ("Please review sections 2 and 4 and approve one option by 17:00 Thursday.", "paragraph"),
            ("Supplier certification remains incomplete; launch moves from 6 May to 20 May.", "paragraph"),
        ]
        blocks = [
            {"blockId": f"burden-{i}", "locator": f"burden:{i}", "text": text, "sourceHash": str(i), "scope": f"burden-{i}", "role": role}
            for i, (text, role) in enumerate(workslop + controls)
        ]
        _, assessments = self.analyzer.analyze(blocks, "general")
        self.assertTrue(all(item["verdict"] == "needs-improvement" for item in assessments[:len(workslop)]))
        self.assertTrue(all(item["verdict"] == "meaningful" for item in assessments[len(workslop):]))

    def test_paraphrased_repetition_is_bounded_by_numbers_and_named_entities(self) -> None:
        repeated = [
            "Customer onboarding takes 14 days because two approvals are manual.",
            "Two manual approvals make customer onboarding last 14 days.",
        ]
        controls = [
            "The target is 8 days after automation.",
            "Revenue increased 14% in Germany.",
            "Revenue increased 14% in France.",
        ]
        blocks = [
            {"blockId": f"repeat-{i}", "locator": f"repeat:{i}", "text": text, "sourceHash": str(i), "scope": "one-slide", "role": "paragraph"}
            for i, text in enumerate(repeated + controls)
        ]
        _, assessments = self.analyzer.analyze(blocks, "consulting")
        self.assertEqual("meaningful", assessments[0]["verdict"])
        self.assertEqual("needs-improvement", assessments[1]["verdict"])
        self.assertTrue(all(item["verdict"] == "meaningful" for item in assessments[2:]))

    def test_unsupported_magnitude_requires_measure_or_source(self) -> None:
        unsupported = [
            "The initiative delivered significant impact across the organization.",
            "Customer satisfaction improved dramatically after implementation.",
            "The pilot was highly successful.",
            "We achieved a substantial reduction in cycle time.",
            "The results were robust and meaningful.",
            "Engagement increased significantly.",
        ]
        controls = [
            "Cycle time fell from 9 days to 4 days.",
            "The pilot met all 12 acceptance criteria.",
            "Sales rose after the price cut.",
            "The difference was significant (Smith, 2024).",
            "The storm intensified rapidly after landfall.",
        ]
        blocks = [
            {"blockId": f"magnitude-{i}", "locator": f"magnitude:{i}", "text": text, "sourceHash": str(i), "scope": f"magnitude-{i}", "role": "paragraph"}
            for i, text in enumerate(unsupported + controls)
        ]
        _, assessments = self.analyzer.analyze(blocks, "general")
        self.assertTrue(all(item["verdict"] == "needs-improvement" for item in assessments[:len(unsupported)]))
        self.assertTrue(all(item["verdict"] == "meaningful" for item in assessments[len(unsupported):]))

    def test_abstract_noun_stacks_fail_without_action_but_keep_concrete_labels(self) -> None:
        empty_stacks = [
            ("Strategic Transformation Enablement", "consulting"),
            ("Enterprise Capability Optimization", "consulting"),
            ("Holistic Innovation Acceleration", "social"),
            ("Operational Excellence Realization", "chat"),
            ("Stakeholder Alignment Facilitation", "chat"),
            ("Future-Ready Value Creation", "social"),
            ("End-to-End Ecosystem Orchestration", "academic"),
            ("Scalable Impact Delivery", "general"),
            ("Organizational Agility Enhancement", "consulting"),
            ("Integrated Solution Implementation", "general"),
        ]
        controls = [
            ("Supplier certification deadline: 31 March", "consulting"),
            ("Revenue forecast for Germany: EUR 4.2m", "consulting"),
            ("Network security audit starts Monday", "general"),
            ("Customer onboarding time fell to 4 days", "social"),
            ("Warehouse fire inspection report", "general"),
            ("Contract approval owner: Finance", "chat"),
            ("Enterprise software implementation starts Monday", "academic"),
            ("Operational excellence team cut defects 12%", "consulting"),
            ("Harvard Innovation Ecosystem includes 42 labs", "academic"),
        ]
        assessments = []
        findings = []
        for i, (text, genre) in enumerate(empty_stacks + controls):
            block = {"blockId": f"noun-stack-{i}", "locator": f"noun-stack:{i}", "text": text, "sourceHash": str(i), "scope": f"noun-stack-{i}", "role": "headline"}
            block_findings, block_assessments = self.analyzer.analyze([block], genre)
            findings.extend(block_findings)
            assessments.extend(block_assessments)
        self.assertTrue(all(item["verdict"] == "needs-improvement" for item in assessments[:len(empty_stacks)]))
        self.assertTrue(all(item["verdict"] == "meaningful" for item in assessments[len(empty_stacks):]))
        noun_findings = [item for item in findings if item["ruleId"] == "VALUE_NOUN_STACK"]
        self.assertEqual(10, len(noun_findings))

    def test_orphan_references_fail_but_clear_container_antecedents_survive(self) -> None:
        orphans = [
            ("This changes everything.", "headline", "social"),
            ("These figures define the opportunity.", "headline", "consulting"),
            ("That is the path forward.", "headline", "consulting"),
            ("It is now clear.", "headline", "academic"),
            ("This is what matters.", "paragraph", "social"),
            ("They need to act now.", "bullet", "chat"),
            ("This unlocks the next phase.", "headline", "general"),
            ("Those are the key priorities.", "paragraph", "consulting"),
        ]
        findings = []
        assessments = []
        for i, (text, role, genre) in enumerate(orphans):
            block = {"blockId": f"orphan-{i}", "locator": f"orphan:{i}", "text": text, "sourceHash": str(i), "scope": f"orphan-{i}", "role": role}
            block_findings, block_assessments = self.analyzer.analyze([block], genre)
            findings.extend(block_findings)
            assessments.extend(block_assessments)
        self.assertTrue(all(item["verdict"] == "needs-improvement" for item in assessments))
        self.assertEqual(8, len([item for item in findings if item["ruleId"] == "VALUE_ORPHAN_REFERENCE"]))

        reordered_slide = [
            {"blockId": "reordered-body", "locator": "reordered:body", "text": "Automation removes two approval steps.", "sourceHash": "rh1", "scope": "reordered", "role": "paragraph"},
            {"blockId": "reordered-headline", "locator": "reordered:headline", "text": "This cuts onboarding time by six days.", "sourceHash": "rh2", "scope": "reordered", "role": "headline"},
        ]
        reordered_findings, reordered_assessments = self.analyzer.analyze(reordered_slide, "consulting")
        self.assertEqual("needs-improvement", reordered_assessments[1]["verdict"])
        self.assertTrue(any(item["ruleId"] == "VALUE_ORPHAN_REFERENCE" and item["blockId"] == "reordered-headline" for item in reordered_findings))

        resolved_pairs = [
            ("The regulator approved the license.", "This allows sales to start Monday.", "consulting"),
            ("Revenue fell 12%.", "This requires a EUR 2m cost reduction.", "consulting"),
            ("The ethics board met on 12 May.", "They approved the protocol that day.", "academic"),
            ("The proposed plant has six production lines.", "It costs EUR 4.2m and takes six months.", "general"),
        ]
        for i, (antecedent, reference, genre) in enumerate(resolved_pairs):
            scope = f"resolved-{i}"
            blocks = [
                {"blockId": f"{scope}-a", "locator": f"{scope}:a", "text": antecedent, "sourceHash": f"{i}a", "scope": scope, "role": "paragraph"},
                {"blockId": f"{scope}-b", "locator": f"{scope}:b", "text": reference, "sourceHash": f"{i}b", "scope": scope, "role": "paragraph"},
            ]
            block_findings, block_assessments = self.analyzer.analyze(blocks, genre)
            self.assertEqual("meaningful", block_assessments[1]["verdict"])
            self.assertFalse(any(item["ruleId"] == "VALUE_ORPHAN_REFERENCE" for item in block_findings))

        self_defining = [
            "These three suppliers failed certification: Alpha, Beta, and Gamma.",
            "This study estimates treatment effects across 42 hospitals.",
            "Those two options cost EUR 2m and EUR 3m, respectively.",
        ]
        for i, text in enumerate(self_defining):
            block = {"blockId": f"defined-{i}", "locator": f"defined:{i}", "text": text, "sourceHash": str(i), "scope": f"defined-{i}", "role": "paragraph"}
            block_findings, block_assessments = self.analyzer.analyze([block], "general")
            self.assertEqual("meaningful", block_assessments[0]["verdict"])
            self.assertFalse(any(item["ruleId"] == "VALUE_ORPHAN_REFERENCE" for item in block_findings))

    def test_explicit_container_contradictions_bind_to_headline(self) -> None:
        conflicts = [
            ("Revenue grew 20% in Germany.", "Revenue fell 8% in Germany."),
            ("The pilot is approved.", "Pilot approval remains pending."),
            ("All suppliers are certified.", "Two suppliers remain uncertified."),
            ("Option A meets all safety requirements.", "Option A fails the safety requirement."),
            ("The launch remains on track for June.", "The June launch has been postponed."),
        ]
        controls = [
            ("Revenue fell in France but grew in Germany.", "France declined 5%; Germany increased 8%."),
            ("The pilot is approved.", "Risk: approval could be withdrawn if the safety test fails."),
            ("Option A meets cost and speed requirements.", "Option A does not meet the optional color preference."),
            ("All production suppliers are certified.", "Two prospective suppliers remain uncertified."),
            ("The launch remains on track for June.", "The June launch may be postponed if certification slips."),
        ]
        blocks = []
        for group, pairs in (("conflict", conflicts), ("control", controls)):
            for index, (headline, body) in enumerate(pairs):
                scope = f"{group}-{index}"
                blocks.extend([
                    {"blockId": f"{scope}-headline", "locator": f"{scope}:headline", "text": headline, "sourceHash": f"{scope}-h", "scope": scope, "role": "headline"},
                    {"blockId": f"{scope}-body", "locator": f"{scope}:body", "text": body, "sourceHash": f"{scope}-b", "scope": scope, "role": "paragraph"},
                ])
        findings, assessments = self.analyzer.analyze(blocks, "consulting")
        by_id = {item["blockId"]: item for item in assessments}
        contradiction_findings = [item for item in findings if item["ruleId"] == "CONTAINER_EXPLICIT_CONTRADICTION"]
        self.assertEqual(5, len(contradiction_findings))
        for index in range(len(conflicts)):
            self.assertEqual("needs-improvement", by_id[f"conflict-{index}-headline"]["verdict"])
        for index in range(len(controls)):
            self.assertEqual("meaningful", by_id[f"control-{index}-headline"]["verdict"])
            self.assertFalse(any(item["blockId"] == f"control-{index}-headline" for item in contradiction_findings))
        self.assertGreater(editorial_vector(contradiction_findings, assessments)["structure"], 0)

    def test_format_adapters_group_headlines_with_supporting_content(self) -> None:
        from docx import Document
        from pptx import Presentation

        with tempfile.TemporaryDirectory() as tmp:
            root = Path(tmp)
            markdown = root / "conflict.md"
            markdown.write_text("# The pilot is approved.\n\nPilot approval remains pending.\n", encoding="utf-8")

            word = root / "conflict.docx"
            document = Document()
            document.add_heading("The pilot is approved.", level=1)
            document.add_paragraph("Pilot approval remains pending.")
            document.save(word)

            powerpoint = root / "conflict.pptx"
            presentation = Presentation()
            slide = presentation.slides.add_slide(presentation.slide_layouts[1])
            slide.shapes.title.text = "The pilot is approved."
            slide.placeholders[1].text = "Pilot approval remains pending."
            presentation.save(powerpoint)

            for path in (markdown, word, powerpoint):
                source = extract_source(request({"kind": "file", "path": str(path)}))
                findings, _ = self.analyzer.analyze(source["blocks"], "consulting")
                conflicts = [item for item in findings if item["ruleId"] == "CONTAINER_EXPLICIT_CONTRADICTION"]
                self.assertEqual(1, len(conflicts), path.suffix)

    def test_every_block_value_coverage(self) -> None:
        source = extract_source(request({"kind": "file", "path": str(self.fixtures / "sample.pptx")}))
        blocks = [block for block in source["blocks"] if block["eligible"]]
        _, assessments = self.analyzer.analyze(blocks, "consulting")
        self.assertEqual(len(blocks), len(assessments))
        self.assertEqual({b["blockId"] for b in blocks}, {a["blockId"] for a in assessments})
        self.assertTrue(any(a["verdict"] == "needs-improvement" for a in assessments))

    def test_protected_tokens(self) -> None:
        before = 'Revenue was 14% (Smith, 2024); see https://example.com and "keep this quote".'
        after = 'Revenue reached 14% (Smith, 2024); see https://example.com and "keep this quote".'
        self.assertTrue(protected_equal(before, after, ["Revenue"])[0])
        self.assertFalse(protected_equal(before, after.replace("14%", "15%"), ["Revenue"])[0])

    def test_strict_request_rejects_unknown_and_duplicates(self) -> None:
        data = request({"kind": "text", "text": "Hello"})
        data["threshold"] = 5
        with self.assertRaises(DeslopError):
            validate_request(data)
        with tempfile.TemporaryDirectory() as td:
            path = Path(td) / "duplicate.json"; path.write_text('{"a": 1, "a": 2}', encoding="utf-8")
            with self.assertRaises(DeslopError):
                strict_json_load(path)

    def test_text_guarded_run_and_independent_verify(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            temp = Path(td); req = request({"kind": "file", "path": str(self.fixtures / "sample.txt")}, "audit-and-rewrite")
            request_path = temp / "request.json"; request_path.write_text(json.dumps(req), encoding="utf-8")
            original = sha256_file(self.fixtures / "sample.txt")
            run = temp / "run"
            call = subprocess.run([sys.executable, str(SKILL / "scripts" / "deslop.py"), "request", str(request_path), "--out", str(run)], capture_output=True, text=True)
            self.assertEqual(0, call.returncode, call.stderr)
            self.assertEqual(original, sha256_file(self.fixtures / "sample.txt"))
            receipt = json.loads((run / "verification.json").read_text(encoding="utf-8"))
            self.assertEqual(receipt["eligibleBlockCount"], receipt["assessedBlockCount"])
            self.assertTrue((run / "value-coverage.json").exists())
            self.assertTrue((run / "sample-deslopped.txt").exists())
            external = temp / "external.json"
            verify = subprocess.run([sys.executable, str(SKILL / "scripts" / "deslop.py"), "request-verify", str(run), "--out", str(external)], capture_output=True, text=True)
            self.assertEqual(0, verify.returncode, verify.stderr)

    def test_hash_bound_semantic_replacement(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            temp = Path(td)
            req = request({"kind": "text", "text": "Mission aligned. Evidence first."}, "audit-and-rewrite")
            req["id"] = "semantic-replacement-test"
            request_path = temp / "request.json"; request_path.write_text(json.dumps(req), encoding="utf-8")
            packet_path = temp / "packet.json"
            ingest = subprocess.run([sys.executable, str(SKILL / "scripts" / "deslop.py"), "ingest", str(request_path), "--out", str(packet_path)], capture_output=True, text=True)
            self.assertEqual(0, ingest.returncode, ingest.stderr)
            packet = json.loads(packet_path.read_text(encoding="utf-8")); block = packet["blocks"][0]
            assessment = {
                "blockId": block["blockId"], "sourceHash": block["sourceHash"], "verdict": "needs-improvement",
                "meaning": "Signals a preference for evidence.", "valueAdded": "Too generic to guide a decision.",
                "relevance": "Prominent copy must state the action.", "reason": "Universal-fit slogan.",
                "improvement": "State the action directly.", "replacement": "Validate one buyer decision first.",
            }
            semantic_path = temp / "semantic.json"; semantic_path.write_text(json.dumps({"schemaVersion": "deslop-semantic/v1", "requestId": req["id"], "assessments": [assessment]}), encoding="utf-8")
            run = temp / "run"
            call = subprocess.run([sys.executable, str(SKILL / "scripts" / "deslop.py"), "request", str(request_path), "--semantic-findings", str(semantic_path), "--out", str(run)], capture_output=True, text=True)
            self.assertEqual(0, call.returncode, call.stderr)
            self.assertEqual("Validate one buyer decision first.", (run / "revised-text.txt").read_text(encoding="utf-8"))
            plan = json.loads((run / "rewrite-plan.json").read_text(encoding="utf-8"))
            self.assertEqual("agent-semantic-source-bound", plan["edits"][0]["rewriteRules"][0])

    def test_format_extractors_and_pdf_audit_policy(self) -> None:
        for name in ["sample.docx", "sample.pptx", "sample.pdf", "sample.md"]:
            source = extract_source(request({"kind": "file", "path": str(self.fixtures / name)}))
            self.assertIn("blocks", source)
            self.assertIn("invariants", source)
        pdf = extract_source(request({"kind": "file", "path": str(self.fixtures / "sample.pdf")}))
        self.assertEqual("pdf", pdf["kind"])
        self.assertTrue(all(not b["supportedForRewrite"] for b in pdf["blocks"]))

    def test_markdown_context_and_numeric_range_are_not_artifacts(self) -> None:
        source = extract_source(request({"kind": "file", "path": str(self.fixtures / "sample.md")}))
        findings, _ = self.analyzer.analyze(source["blocks"], "general")
        ids = {item["ruleId"] for item in findings}
        self.assertNotIn("ART_MARKDOWN_HEADING", ids)
        self.assertNotIn("STR_RANGE_FROM_TO", ids)

    def test_docx_and_pptx_guarded_revisions_open_in_office(self) -> None:
        with tempfile.TemporaryDirectory() as td:
            temp = Path(td)
            for name in ["sample.docx", "sample.pptx"]:
                req = request({"kind": "file", "path": str(self.fixtures / name)}, "audit-and-rewrite")
                req["id"] = f"synthetic-{Path(name).stem}"
                request_path = temp / f"{name}.json"; request_path.write_text(json.dumps(req), encoding="utf-8")
                run = temp / f"run-{Path(name).suffix.lstrip('.')}"
                original = sha256_file(self.fixtures / name)
                call = subprocess.run([sys.executable, str(SKILL / "scripts" / "deslop.py"), "request", str(request_path), "--out", str(run)], capture_output=True, text=True, timeout=90)
                self.assertEqual(0, call.returncode, f"{name}: {call.stderr}\n{call.stdout}")
                self.assertEqual(original, sha256_file(self.fixtures / name))
                receipt = json.loads((run / "verification.json").read_text(encoding="utf-8"))
                self.assertTrue(receipt["passed"])
                self.assertTrue((run / f"sample-deslopped{Path(name).suffix}").exists())
                plan = json.loads((run / "rewrite-plan.json").read_text(encoding="utf-8"))
                self.assertTrue(any("mixed run formatting" in item["reason"].casefold() for item in plan["abstentions"]))

    def test_private_profile_stores_metrics_not_text(self) -> None:
        PRIVATE_ROOT.mkdir(parents=True, exist_ok=True)
        with tempfile.TemporaryDirectory(dir=PRIVATE_ROOT) as td:
            parent = Path(td); profile_request = {
                "schemaVersion": "deslop-profile-request/v1", "id": "consulting-test", "genre": "consulting",
                "inputs": [str(self.fixtures)], "approvedExcerptIds": [],
            }
            request_path = parent / "profile-request-input.json"; request_path.write_text(json.dumps(profile_request), encoding="utf-8")
            output = parent / "profile"
            call = subprocess.run([sys.executable, str(SKILL / "scripts" / "deslop.py"), "profile", str(request_path), "--out", str(output)], capture_output=True, text=True)
            self.assertEqual(0, call.returncode, call.stderr)
            profile = json.loads((output / "profile.json").read_text(encoding="utf-8"))
            self.assertFalse(profile["privacy"]["rawTextStored"])
            serialized = json.dumps(profile)
            self.assertNotIn("Unlocking Transformative Excellence", serialized)
            self.assertNotIn("documents", profile["corpus"])
            self.assertNotIn("@", serialized)
            self.assertNotIn(str(self.fixtures), serialized)
            self.assertFalse((output / "profile-request.json").exists())


if __name__ == "__main__":
    unittest.main()
