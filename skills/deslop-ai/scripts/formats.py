from __future__ import annotations

import re
import zipfile
from pathlib import Path
from typing import Any, Iterable

from deslop_core import DeslopError, sha256_file, sha256_text, stable_id


def _package_hashes(path: Path, prefixes: tuple[str, ...], suffixes: tuple[str, ...] = ()) -> dict[str, str]:
    hashes: dict[str, str] = {}
    with zipfile.ZipFile(path) as package:
        for name in sorted(package.namelist()):
            if name.startswith(prefixes) or (suffixes and name.endswith(suffixes)):
                hashes[name] = sha256_text(package.read(name).hex())
    return hashes


def _relationship_inventory(path: Path) -> dict[str, list[list[str]]]:
    from lxml import etree
    inventory: dict[str, list[list[str]]] = {}
    with zipfile.ZipFile(path) as package:
        for name in sorted(item for item in package.namelist() if item.endswith(".rels")):
            root = etree.fromstring(package.read(name))
            entries = []
            for relationship in root:
                entries.append(sorted([
                    relationship.get("Type", ""),
                    relationship.get("Target", ""),
                    relationship.get("TargetMode", ""),
                ]))
            inventory[name] = sorted(entries)
    return inventory


def _auxiliary_word_blocks(path: Path) -> list[dict[str, Any]]:
    blocks: list[dict[str, Any]] = []
    try:
        from lxml import etree
        with zipfile.ZipFile(path) as package:
            for part in ("word/footnotes.xml", "word/endnotes.xml", "word/comments.xml"):
                if part not in package.namelist():
                    continue
                root = etree.fromstring(package.read(part))
                namespaces = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
                for index, paragraph in enumerate(root.xpath(".//w:p", namespaces=namespaces)):
                    text = "".join(paragraph.xpath(".//w:t/text()", namespaces=namespaces)).strip()
                    if text:
                        label = part.split("/")[-1].split(".")[0]
                        blocks.append(_block(
                            f"{label}:paragraph:{index}", text, scope=label, role="notes", format="docx",
                            address={"part": label, "paragraph": index}, supportedForRewrite=False,
                            unsupportedReason=f"{label} are audit-only in v1",
                        ))
            if "word/document.xml" in package.namelist():
                root = etree.fromstring(package.read("word/document.xml"))
                namespaces = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
                for label, expression in (("content-control", ".//w:sdt"), ("tracked-insertion", ".//w:ins"), ("tracked-deletion", ".//w:del")):
                    for index, node in enumerate(root.xpath(expression, namespaces=namespaces)):
                        text = "".join(node.xpath(".//w:t/text()", namespaces=namespaces)).strip()
                        if text:
                            blocks.append(_block(
                                f"{label}:{index}", text, scope=label, role="notes", format="docx",
                                address={"part": label, "index": index}, supportedForRewrite=False,
                                unsupportedReason=f"{label} text is audit-only in v1",
                            ))
    except (KeyError, zipfile.BadZipFile):
        pass
    return blocks


def _block(locator: str, text: str, **extra: Any) -> dict[str, Any]:
    return {
        "blockId": stable_id("block", locator, sha256_text(text)),
        "locator": locator,
        "text": text,
        "sourceHash": sha256_text(text),
        "eligible": bool(text.strip()),
        "supportedForRewrite": True,
        **extra,
    }


def _paragraph_role(text: str, style: str = "", is_bullet: bool = False) -> str:
    if re.fullmatch(r"\s*\d+\s*", text):
        return "furniture"
    if re.match(r"\s*(?:confidential|draft|source|sources|©|copyright)\b", text, re.I):
        return "furniture"
    if "title" in style.casefold() or "heading" in style.casefold():
        return "headline"
    if is_bullet:
        return "bullet"
    return "paragraph"


def extract_text_file(path: Path) -> dict[str, Any]:
    raw = path.read_bytes()
    encoding = "utf-8-sig" if raw.startswith(b"\xef\xbb\xbf") else "utf-8"
    try:
        text = raw.decode(encoding)
    except UnicodeDecodeError:
        encoding = "cp1252"
        text = raw.decode(encoding)
    blocks = []
    position = 0
    fenced = False
    for line_index, line_with_end in enumerate(text.splitlines(keepends=True)):
        line = line_with_end.rstrip("\r\n")
        stripped = line.strip()
        if path.suffix.lower() == ".md" and stripped.startswith(chr(96) * 3):
            fenced = not fenced
        start = position
        end = position + len(line)
        position += len(line_with_end)
        if not stripped:
            continue
        is_heading = path.suffix.lower() == ".md" and bool(re.match(r"^\s*#{1,6}\s+", line))
        is_bullet = bool(re.match(r"^\s*[-*+]\s+", line))
        role = "headline" if is_heading else ("bullet" if is_bullet else "paragraph")
        unsupported = fenced or (path.suffix.lower() == ".md" and (chr(96) in line or re.search(r"\[[^]]+\]\([^)]+\)", line)))
        blocks.append(_block(
            f"line:{line_index + 1}", line, scope="document",
            role=role, isBullet=is_bullet, level=0 if is_heading else None,
            format=path.suffix.lower().lstrip("."), address={"start": start, "end": end},
            supportedForRewrite=not unsupported,
            unsupportedReason="Markdown code/link content" if unsupported else "",
        ))
    return {
        "kind": path.suffix.lower().lstrip("."), "path": str(path), "sourceHash": sha256_file(path),
        "encoding": encoding, "newline": "\r\n" if b"\r\n" in raw else "\n", "blocks": blocks,
        "invariants": {"byteLength": len(raw), "lineCount": len(text.splitlines())},
    }


def extract_pasted_text(text: str) -> dict[str, Any]:
    blocks = []
    position = 0
    for index, part in enumerate(re.split(r"(\n\s*\n)", text)):
        start = position
        position += len(part)
        if not part.strip() or re.fullmatch(r"\n\s*\n", part):
            continue
        blocks.append(_block(
            f"text:block:{index}", part, scope="pasted-text", role="paragraph",
            format="text", address={"start": start, "end": position},
        ))
    return {"kind": "text", "path": None, "sourceHash": sha256_text(text), "text": text, "blocks": blocks, "invariants": {"characterCount": len(text)}}


def _docx_run_signature(run: Any) -> tuple[Any, ...]:
    font = run.font
    color = None
    try:
        color = str(font.color.rgb) if font.color and font.color.rgb else None
    except Exception:
        color = None
    return (run.style.name if run.style else None, font.name, font.size.pt if font.size else None, font.bold, font.italic, font.underline, color)


def _docx_supported(paragraph: Any) -> tuple[bool, str]:
    xml = paragraph._p.xml
    forbidden = ("<w:hyperlink", "<w:fldChar", "<w:instrText", "<m:oMath", "<w:sdt", "<w:ins", "<w:del")
    if any(token in xml for token in forbidden):
        return False, "Paragraph contains hyperlink, field, equation, content control, or tracked change"
    runs = [run for run in paragraph.runs if run.text]
    if len({_docx_run_signature(run) for run in runs}) > 1:
        return False, "Paragraph contains mixed run formatting"
    return True, ""


def extract_docx(path: Path) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError as exc:
        raise DeslopError("python-docx is required for DOCX support") from exc
    doc = Document(path)
    blocks: list[dict[str, Any]] = []
    for i, paragraph in enumerate(doc.paragraphs):
        if not paragraph.text.strip():
            continue
        supported, reason = _docx_supported(paragraph)
        style = paragraph.style.name if paragraph.style else ""
        is_bullet = "list" in style.casefold() or bool(paragraph._p.xpath("./w:pPr/w:numPr"))
        blocks.append(_block(
            f"body:paragraph:{i}", paragraph.text, scope=f"body:{i}", role=_paragraph_role(paragraph.text, style, is_bullet),
            style=style, isBullet=is_bullet, level=0 if "title" in style.casefold() else None,
            format="docx", address={"part": "body", "paragraph": i}, supportedForRewrite=supported, unsupportedReason=reason,
        ))
    for ti, table in enumerate(doc.tables):
        for ri, row in enumerate(table.rows):
            for ci, cell in enumerate(row.cells):
                for pi, paragraph in enumerate(cell.paragraphs):
                    if not paragraph.text.strip():
                        continue
                    supported, reason = _docx_supported(paragraph)
                    blocks.append(_block(
                        f"table:{ti}:row:{ri}:cell:{ci}:paragraph:{pi}", paragraph.text, scope=f"table:{ti}:row:{ri}",
                        role="table-cell", container="table", style=paragraph.style.name if paragraph.style else "",
                        isBullet=False, format="docx", address={"part": "table", "table": ti, "row": ri, "cell": ci, "paragraph": pi},
                        supportedForRewrite=supported, unsupportedReason=reason,
                    ))
    for si, section in enumerate(doc.sections):
        for part_name, part in (("header", section.header), ("footer", section.footer)):
            for pi, paragraph in enumerate(part.paragraphs):
                if paragraph.text.strip():
                    blocks.append(_block(
                        f"section:{si}:{part_name}:paragraph:{pi}", paragraph.text, scope=f"section:{si}:{part_name}",
                        role="furniture", format="docx", address={"part": part_name, "section": si, "paragraph": pi},
                        supportedForRewrite=False, unsupportedReason=f"{part_name} is audit-only in v1",
                    ))
    blocks.extend(_auxiliary_word_blocks(path))
    with zipfile.ZipFile(path) as package:
        xml = b"\n".join(package.read(name) for name in package.namelist() if name.endswith(".xml"))
    invariants = {
        "paragraphCount": len(doc.paragraphs), "tableCount": len(doc.tables), "sectionCount": len(doc.sections),
        "tableDimensions": [[len(table.rows), len(table.columns)] for table in doc.tables],
        "styleNames": sorted(style.name for style in doc.styles),
        "relationships": _relationship_inventory(path),
        "nonTextParts": _package_hashes(path, ("word/media/", "word/embeddings/")),
        "fieldCount": xml.count(b"fldChar") + xml.count(b"instrText"),
        "contentControlCount": xml.count(b"<w:sdt"),
        "trackedChangeCount": xml.count(b"<w:ins") + xml.count(b"<w:del"),
    }
    return {"kind": "docx", "path": str(path), "sourceHash": sha256_file(path), "blocks": blocks, "invariants": invariants}


def _pptx_run_signature(run: Any) -> tuple[Any, ...]:
    font = run.font
    color = None
    try:
        color = str(font.color.rgb) if font.color and font.color.type else None
    except Exception:
        color = None
    hyperlink = None
    try:
        hyperlink = run.hyperlink.address
    except Exception:
        hyperlink = None
    return (font.name, font.size.pt if font.size else None, font.bold, font.italic, font.underline, color, hyperlink)


def _pptx_supported(paragraph: Any) -> tuple[bool, str]:
    runs = [run for run in paragraph.runs if run.text]
    signatures = {_pptx_run_signature(run) for run in runs}
    if any(signature[-1] for signature in signatures):
        return False, "Paragraph contains a hyperlink"
    if len(signatures) > 1:
        return False, "Paragraph contains mixed run formatting"
    if "fld" in paragraph._p.xml:
        return False, "Paragraph contains a field"
    return True, ""


def _walk_shapes(shapes: Any, prefix: tuple[int, ...] = ()) -> Iterable[tuple[tuple[int, ...], Any]]:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    for index, shape in enumerate(shapes):
        path = prefix + (index,)
        yield path, shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _walk_shapes(shape.shapes, path)


def _shape_by_path(shapes: Any, path: list[int]) -> Any:
    current = shapes
    shape = None
    for index in path:
        shape = current[index]
        current = shape.shapes if hasattr(shape, "shapes") else None
    return shape


def extract_pptx(path: Path) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DeslopError("python-pptx is required for PPTX support") from exc
    prs = Presentation(path)
    blocks: list[dict[str, Any]] = []
    geometry = []
    for si, slide in enumerate(prs.slides):
        for shape_path, shape in _walk_shapes(slide.shapes):
            shape_id = getattr(shape, "shape_id", None)
            geometry.append({
                "slide": si, "path": list(shape_path), "shapeId": shape_id, "type": str(shape.shape_type),
                "left": int(shape.left), "top": int(shape.top), "width": int(shape.width), "height": int(shape.height),
            })
            scope = f"slide:{si + 1}"
            if getattr(shape, "has_text_frame", False):
                for pi, paragraph in enumerate(shape.text_frame.paragraphs):
                    if not paragraph.text.strip():
                        continue
                    supported, reason = _pptx_supported(paragraph)
                    placeholder_type = getattr(getattr(shape, "placeholder_format", None), "type", None) if getattr(shape, "is_placeholder", False) else None
                    is_title = bool(placeholder_type in (1, 3))
                    is_furniture = bool(placeholder_type in (10, 11, 12, 13))
                    paragraph_xml = paragraph._p.xml
                    is_body_placeholder = bool(placeholder_type in (2, 7, 14, 15, 16, 17, 18))
                    is_bullet = paragraph.level > 0 or "<a:buChar" in paragraph_xml or "<a:buAutoNum" in paragraph_xml or is_body_placeholder
                    role = "furniture" if is_furniture else ("headline" if is_title else ("bullet" if is_bullet else ("headline" if shape.top < prs.slide_height * 0.25 and len(paragraph.text) <= 120 else "paragraph")))
                    blocks.append(_block(
                        f"slide:{si + 1}:shape:{'.'.join(map(str, shape_path))}:paragraph:{pi}", paragraph.text, scope=scope,
                        role=role, isBullet=role == "bullet", level=int(paragraph.level), format="pptx",
                        address={"slide": si, "shapePath": list(shape_path), "paragraph": pi, "container": "shape"},
                        supportedForRewrite=supported, unsupportedReason=reason,
                    ))
            if getattr(shape, "has_table", False):
                for ri, row in enumerate(shape.table.rows):
                    for ci, cell in enumerate(row.cells):
                        for pi, paragraph in enumerate(cell.text_frame.paragraphs):
                            if not paragraph.text.strip():
                                continue
                            supported, reason = _pptx_supported(paragraph)
                            blocks.append(_block(
                                f"slide:{si + 1}:shape:{'.'.join(map(str, shape_path))}:table:{ri}:{ci}:paragraph:{pi}",
                                paragraph.text, scope=scope, role="table-cell", container="table", isBullet=False, level=int(paragraph.level),
                                format="pptx", address={"slide": si, "shapePath": list(shape_path), "paragraph": pi, "container": "table", "row": ri, "cell": ci},
                                supportedForRewrite=supported, unsupportedReason=reason,
                            ))
            if getattr(shape, "has_chart", False):
                chart = shape.chart
                if chart.has_title:
                    for pi, paragraph in enumerate(chart.chart_title.text_frame.paragraphs):
                        if paragraph.text.strip():
                            blocks.append(_block(
                                f"slide:{si + 1}:shape:{'.'.join(map(str, shape_path))}:chart:title:{pi}", paragraph.text,
                                scope=scope, role="caption", format="pptx", address={"slide": si, "shapePath": list(shape_path), "container": "chart-title", "paragraph": pi},
                                supportedForRewrite=False, unsupportedReason="Chart text is audit-only in v1",
                            ))
                for series_index, series in enumerate(chart.series):
                    if str(series.name).strip():
                        blocks.append(_block(
                            f"slide:{si + 1}:shape:{'.'.join(map(str, shape_path))}:chart:series:{series_index}", str(series.name),
                            scope=scope, role="chart-label", format="pptx", address={"slide": si, "shapePath": list(shape_path), "container": "chart-series", "series": series_index},
                            supportedForRewrite=False, unsupportedReason="Chart text is audit-only in v1",
                        ))
                try:
                    categories = list(chart.plots[0].categories)
                    for category_index, category in enumerate(categories):
                        label = str(getattr(category, "label", category)).strip()
                        if label:
                            blocks.append(_block(
                                f"slide:{si + 1}:shape:{'.'.join(map(str, shape_path))}:chart:category:{category_index}", label,
                                scope=scope, role="chart-label", format="pptx", address={"slide": si, "shapePath": list(shape_path), "container": "chart-category", "category": category_index},
                                supportedForRewrite=False, unsupportedReason="Chart text is audit-only in v1",
                            ))
                except Exception:
                    pass
        try:
            notes = slide.notes_slide.notes_text_frame
            for pi, paragraph in enumerate(notes.paragraphs):
                if paragraph.text.strip():
                    blocks.append(_block(
                        f"slide:{si + 1}:notes:paragraph:{pi}", paragraph.text, scope=f"slide:{si + 1}:notes", role="notes",
                        format="pptx", address={"slide": si, "paragraph": pi, "container": "notes"},
                        supportedForRewrite=False, unsupportedReason="Speaker notes are audit-only in v1",
                    ))
        except Exception:
            pass
    try:
        from lxml import etree
        with zipfile.ZipFile(path) as package:
            diagram_parts = sorted(name for name in package.namelist() if name.startswith("ppt/diagrams/data") and name.endswith(".xml"))
            for part_index, part in enumerate(diagram_parts):
                root = etree.fromstring(package.read(part))
                texts = root.xpath(".//*[local-name()='t']/text()")
                for text_index, text in enumerate(texts):
                    if str(text).strip():
                        blocks.append(_block(
                            f"smartart:part:{part_index}:text:{text_index}", str(text).strip(), scope=f"smartart:{part_index}",
                            role="smartart-text", format="pptx", address={"container": "smartart", "part": part, "text": text_index},
                            supportedForRewrite=False, unsupportedReason="SmartArt text is audit-only in v1",
                        ))
    except (KeyError, zipfile.BadZipFile):
        pass
    invariants = {
        "slideCount": len(prs.slides), "slideWidth": int(prs.slide_width), "slideHeight": int(prs.slide_height),
        "masterCount": len(prs.slide_masters), "layoutCount": sum(len(master.slide_layouts) for master in prs.slide_masters),
        "geometry": geometry,
        "slideBindings": [
            {
                "layout": slide.slide_layout.name,
                "master": slide.slide_layout.slide_master.name,
            }
            for slide in prs.slides
        ],
        "nonTextParts": _package_hashes(path, ("ppt/media/", "ppt/charts/", "ppt/diagrams/", "ppt/embeddings/", "ppt/theme/")),
        "relationships": _relationship_inventory(path),
    }
    return {"kind": "pptx", "path": str(path), "sourceHash": sha256_file(path), "blocks": blocks, "invariants": invariants}


def extract_pdf(path: Path) -> dict[str, Any]:
    try:
        from pypdf import PdfReader
    except ImportError as exc:
        raise DeslopError("pypdf is required for PDF support") from exc
    reader = PdfReader(path)
    blocks = []
    for page_index, page in enumerate(reader.pages):
        text = page.extract_text() or ""
        for paragraph_index, paragraph in enumerate(re.split(r"\n\s*\n|\n(?=[A-Z])", text)):
            if paragraph.strip():
                blocks.append(_block(
                    f"page:{page_index + 1}:block:{paragraph_index}", paragraph.strip(), scope=f"page:{page_index + 1}",
                    role="paragraph", format="pdf", address={"page": page_index, "block": paragraph_index},
                    supportedForRewrite=False, unsupportedReason="PDF is audit-only in v1",
                ))
    return {"kind": "pdf", "path": str(path), "sourceHash": sha256_file(path), "blocks": blocks, "invariants": {"pageCount": len(reader.pages)}}


def extract_source(request: dict[str, Any]) -> dict[str, Any]:
    source = request["input"]
    if source["kind"] == "text":
        return extract_pasted_text(source["text"])
    path = Path(source["path"])
    suffix = path.suffix.lower()
    if suffix in {".txt", ".md"}:
        return extract_text_file(path)
    if suffix == ".docx":
        return extract_docx(path)
    if suffix == ".pptx":
        return extract_pptx(path)
    if suffix == ".pdf":
        return extract_pdf(path)
    raise DeslopError(f"Unsupported format: {suffix}")


def _replace_paragraph_runs(paragraph: Any, replacement: str) -> None:
    runs = list(paragraph.runs)
    if runs:
        runs[0].text = replacement
        for run in runs[1:]:
            run.text = ""
    else:
        paragraph.add_run(replacement)


def apply_text_edits(source_map: dict[str, Any], edits: list[dict[str, Any]], output: Path) -> dict[str, Any]:
    if source_map["kind"] == "text":
        text = source_map["text"]
    else:
        source = Path(source_map["path"])
        raw = source.read_bytes()
        encoding = source_map["encoding"]
        text = raw.decode(encoding)
    for edit in sorted(edits, key=lambda item: item["address"]["start"], reverse=True):
        start, end = edit["address"]["start"], edit["address"]["end"]
        if sha256_text(text[start:end]) != edit["oldHash"]:
            raise DeslopError(f"Stale text edit: {edit['editId']}")
        text = text[:start] + edit["replacement"] + text[end:]
    output_encoding = "utf-8" if source_map["kind"] == "text" else source_map["encoding"]
    output.write_bytes(text.encode(output_encoding))
    return {"characterCount": len(text)}


def apply_docx_edits(source_map: dict[str, Any], edits: list[dict[str, Any]], output: Path) -> dict[str, Any]:
    from docx import Document
    source = Path(source_map["path"])
    doc = Document(source)
    for edit in edits:
        address = edit["address"]
        if address["part"] == "body":
            paragraph = doc.paragraphs[address["paragraph"]]
        elif address["part"] == "table":
            paragraph = doc.tables[address["table"]].rows[address["row"]].cells[address["cell"]].paragraphs[address["paragraph"]]
        else:
            raise DeslopError(f"Unsupported DOCX edit target: {address['part']}")
        if sha256_text(paragraph.text) != edit["oldHash"]:
            raise DeslopError(f"Stale DOCX edit: {edit['editId']}")
        supported, reason = _docx_supported(paragraph)
        if not supported:
            raise DeslopError(f"Unsafe DOCX edit {edit['editId']}: {reason}")
        _replace_paragraph_runs(paragraph, edit["replacement"])
    doc.save(output)
    return extract_docx(output)["invariants"]


def apply_pptx_edits(source_map: dict[str, Any], edits: list[dict[str, Any]], output: Path) -> dict[str, Any]:
    from pptx import Presentation
    source = Path(source_map["path"])
    prs = Presentation(source)
    for edit in edits:
        address = edit["address"]
        slide = prs.slides[address["slide"]]
        if address["container"] == "notes":
            raise DeslopError("Notes are audit-only")
        shape = _shape_by_path(slide.shapes, address["shapePath"])
        if address["container"] == "table":
            paragraph = shape.table.rows[address["row"]].cells[address["cell"]].text_frame.paragraphs[address["paragraph"]]
        else:
            paragraph = shape.text_frame.paragraphs[address["paragraph"]]
        if sha256_text(paragraph.text) != edit["oldHash"]:
            raise DeslopError(f"Stale PPTX edit: {edit['editId']}")
        supported, reason = _pptx_supported(paragraph)
        if not supported:
            raise DeslopError(f"Unsafe PPTX edit {edit['editId']}: {reason}")
        _replace_paragraph_runs(paragraph, edit["replacement"])
    prs.save(output)
    return extract_pptx(output)["invariants"]


def apply_edits(source_map: dict[str, Any], edits: list[dict[str, Any]], output: Path) -> dict[str, Any]:
    output.parent.mkdir(parents=True, exist_ok=True)
    kind = source_map["kind"]
    if kind in {"text", "txt", "md"}:
        return apply_text_edits(source_map, edits, output)
    if kind == "docx":
        return apply_docx_edits(source_map, edits, output)
    if kind == "pptx":
        return apply_pptx_edits(source_map, edits, output)
    if kind == "pdf":
        raise DeslopError("PDF rewriting is not supported")
    raise DeslopError(f"Unsupported apply format: {kind}")


def compare_invariants(before: dict[str, Any], after: dict[str, Any], kind: str) -> list[str]:
    problems = []
    if kind == "pptx":
        for key in ("slideCount", "slideWidth", "slideHeight", "masterCount", "layoutCount", "slideBindings", "geometry", "nonTextParts", "relationships"):
            if before.get(key) != after.get(key):
                problems.append(f"PPTX invariant changed: {key}")
    elif kind == "docx":
        for key in ("paragraphCount", "tableCount", "sectionCount", "tableDimensions", "styleNames", "relationships", "nonTextParts", "fieldCount", "contentControlCount", "trackedChangeCount"):
            if before.get(key) != after.get(key):
                problems.append(f"DOCX invariant changed: {key}")
    return problems
